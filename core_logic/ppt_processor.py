# core_logic/ppt_processor.py
import os
import platform
import time
import shutil
from pathlib import Path
import logging
import uuid
import configparser
import json
import subprocess
import shlex
import sys
import traceback # 用于获取堆栈跟踪

# 导入 mutagen (如果需要)
try:
    from mutagen import File as MutagenFile, MutagenError
    MUTAGEN_AVAILABLE = True
except ImportError:
    logging.warning("缺少 'mutagen' 库，MP3 时长可能不准。'pip install mutagen'")
    MUTAGEN_AVAILABLE = False


# 导入同级模块使用相对路径
from .tts_manager_edge import generate_audio_segment # 导入 tts 管理器
from .ppt_exporter_libreoffice import export_slides_with_libreoffice # 只导入 LibreOffice 导出器
from .utils import get_audio_duration, get_tool_path # 从 utils 导入工具函数


# 导入 Presentation 类
try:
    from pptx import Presentation
except ImportError:
    logging.error("FATAL ERROR: 缺少 'python-pptx' 库！请运行 'pip install python-pptx'")
    # 在模块级别无法导入，需要处理。可以在调用此模块的 Celery 任务开始时检查。
    raise ImportError("缺少必需的 python-pptx 库")


# --- 定义任务阶段常量 (用于状态更新) ---
# 这些常量也在 tasks.py 中定义，确保一致
STAGE_START = 'Initializing'
STAGE_PPT_PROCESSING = 'Processing Presentation'
STAGE_PPT_IMAGES = 'Exporting Slides'
STAGE_EXTRACT_NOTES = 'Extracting Notes'
STAGE_GENERATE_AUDIO = 'Generating Audio'
STAGE_VIDEO_SYNTHESIS = 'Synthesizing Video'
STAGE_VIDEO_SEGMENTS = 'Creating Video Segments'
STAGE_VIDEO_CONCAT = 'Concatenating Video'
STAGE_GENERATE_SUBTITLES = 'Generating Subtitles (ASR)'
STAGE_ADD_SUBTITLES = 'Adding Subtitles'
STAGE_CLEANUP = 'Cleaning Up'
STAGE_COMPLETE = 'Complete'


def extract_speaker_notes(pptx_filepath: Path, logger: logging.Logger, task_instance) -> list[str]:
    """
    从 PPTX 文件提取演讲者备注。失败时抛出异常。
    增加 task_instance 参数，但暂时不在此函数内部发送状态。
    """
    if not pptx_filepath.is_file():
        logger.error(f"输入文件不存在: {pptx_filepath}")
        raise FileNotFoundError(f"输入文件不存在: {pptx_filepath}")

    notes_list = []
    try:
        # task_instance.update_state('PROCESSING', meta={'stage': STAGE_EXTRACT_NOTES, 'status': 'Starting'})
        logger.info(f"开始解析演示文稿以提取备注: {pptx_filepath.name}")
        prs = Presentation(pptx_filepath)
        num_slides = len(prs.slides)
        logger.info(f"演示文稿包含 {num_slides} 张幻灯片。")

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            note_text = "" # 默认为空字符串
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame and text_frame.text:
                    note_text = text_frame.text.strip()
                    # logger.debug(f"  找到幻灯片 {slide_num} 的备注: '{note_text[:50]}...'") # 记录备注开头部分
            # task_instance.update_state('PROCESSING', meta={'stage': STAGE_EXTRACT_NOTES, 'progress': int((i + 1) / num_slides * 100), 'current_slide': slide_num})

            notes_list.append(note_text)

        logger.info(f"成功提取了 {len(notes_list)} 条备注信息。")
        # task_instance.update_state('PROCESSING', meta={'stage': STAGE_EXTRACT_NOTES, 'progress': 100, 'status': 'Notes extracted'})
        return notes_list

    except Exception as e:
        logger.error(f"解析 PPTX 文件以提取备注时出错: {e}", exc_info=True)
        # task_instance.update_state('PROCESSING', meta={'stage': STAGE_EXTRACT_NOTES, 'status': f'Error: Extraction failed ({type(e).__name__})'})
        raise Exception(f"提取备注失败: {e}") from e


# --- generate_audio_segments 函数 (修改调用，增加 task_instance) ---
def generate_audio_segments(
    notes: list[str],
    output_audio_dir: Path,
    voice_id: str,
    rate: int, # 语速百分比
    logger: logging.Logger,
    config: configparser.ConfigParser,
    task_instance # <--- 接收任务实例
) -> list[tuple[str | None, float]]: # 返回有效时长 (float)，失败记为 0.0
    """
    使用 Edge TTS 将文本备注列表转换为 MP3 音频文件。

    Args:
        notes: 包含每张幻灯片备注文本的字符串列表。
        output_audio_dir: 保存生成的 MP3 文件的目标目录。
        voice_id: 要使用的 Edge TTS 语音 ID (必需)。
        rate: 语速百分比。
        logger: 日志记录器实例。
        config: ConfigParser 对象。
        task_instance: Celery 任务实例 (用于状态更新)。

    Returns:
        一个元组列表，每个元组包含 (生成的音频文件绝对路径字符串 | None, 音频时长 float)。
        生成失败或时长无效则时长记为 0.0。
    """
    audio_results = []
    output_audio_dir.mkdir(parents=True, exist_ok=True)

    logger.info(f"开始使用 Edge TTS 生成音频片段 (Voice ID: {voice_id})...")
    tts_retries = config.getint('Audio', 'tts_retries', fallback=1)
    tts_retry_delay = config.getfloat('Audio', 'tts_retry_delay', fallback=1.5)
    num_segments = len(notes)

    for i, text in enumerate(notes):
        segment_num = i + 1
        # 使用 pathlib 构建路径
        audio_filename = f"segment_{segment_num}_{uuid.uuid4().hex[:4]}.mp3" # 使用唯一文件名
        audio_filepath = output_audio_dir / audio_filename
        audio_path_str = None
        duration_sec = 0.0 # 默认时长为 0.0

        if text and not text.isspace():
            # task_instance.update_state('PROCESSING', meta={'stage': STAGE_GENERATE_AUDIO, 'progress': int(i / num_segments * 100), 'current_segment': segment_num})
            logger.debug(f"  生成片段 {segment_num} 的音频 (文本: '{text[:50]}...')...")

            # 调用 tts_manager 中的函数
            success = generate_audio_segment( # generate_audio_segment 内部会处理重试和异常
                voice_id,
                text,
                audio_filepath,
                rate=rate,
                logger=logger,
                max_retries=tts_retries,
                retry_delay=tts_retry_delay
            )

            if success:
                # 音频生成成功，尝试获取时长
                # logger, config 已经传递
                duration_sec_raw = get_audio_duration(audio_filepath, logger, config)
                if duration_sec_raw is not None and duration_sec_raw > 0.01:
                    duration_sec = duration_sec_raw
                    audio_path_str = str(audio_filepath.resolve())
                    logger.debug(f"    片段 {segment_num} 生成成功, 时长: {duration_sec:.3f}s")
                elif duration_sec_raw is not None: # 时长为 0 或过短
                     audio_path_str = str(audio_filepath.resolve()) # 文件存在
                     logger.warning(f"    片段 {segment_num} 音频时长无效或过短 ({duration_sec_raw:.3f}s)，时长记为 0。")
                     duration_sec = 0.0
                else: # 获取时长失败
                     # get_audio_duration 内部会记录错误
                     audio_path_str = str(audio_filepath.resolve()) # 文件可能存在
                     duration_sec = 0.0 # 时长记为 0
                     logger.error(f"    无法获取片段 {segment_num} ({audio_filepath.name}) 的有效时长！时长记为 0。")
            else: # TTS 生成失败 (generate_audio_segment 返回 False)
                duration_sec = 0.0 # 时长记为 0
                logger.error(f"    片段 {segment_num} TTS 生成失败。")
        else:
            logger.debug(f"  片段 {segment_num}: 文本为空，跳过音频生成，时长为 0。")
            duration_sec = 0.0 # 空文本，时长为 0

        audio_results.append((audio_path_str, duration_sec)) # 存储结果

    # task_instance.update_state('PROCESSING', meta={'stage': STAGE_GENERATE_AUDIO, 'progress': 100, 'status': 'Audio generation complete'})
    logger.info(f"音频生成过程完成。生成了 {len([p for p, d in audio_results if p])} 个有效音频片段。")
    return audio_results


# --- process_presentation_for_task 函数 (修改后，增加 task_instance) ---
def process_presentation_for_task(
    pptx_filepath: Path,
    temp_base_dir: Path,
    voice_id: str,
    logger: logging.Logger,
    config: configparser.ConfigParser,
    task_instance # <--- 接收任务实例
) -> tuple[list[dict], Path]:
    """
    处理演示文稿的核心后台任务逻辑。
    包括：导出幻灯片图片、提取备注、生成音频片段。
    失败时应抛出异常，成功时返回处理后的数据和临时目录路径。
    """
    task_id = task_instance.request.id
    logger.debug(f"任务 {task_id} 调用 process_presentation_for_task")

    if not pptx_filepath.is_file():
        logger.error(f"输入 PPTX 文件不存在: {pptx_filepath}")
        raise FileNotFoundError(f"输入 PPTX 文件不存在: {pptx_filepath}")
    if not voice_id:
        raise ValueError("必须提供有效的 voice_id!")
    if not temp_base_dir or not isinstance(temp_base_dir, Path):
        raise ValueError("必须提供有效的临时目录基础路径 (temp_base_dir)")

    # --- 创建本次任务的独立临时目录 ---
    run_id = uuid.uuid4().hex[:8]
    safe_stem = "".join(c if c.isalnum() or c in ('-', '_') else '_' for c in pptx_filepath.stem)
    temp_run_dir = temp_base_dir / f"task_{run_id}_{safe_stem}"
    temp_image_dir = temp_run_dir / "images"
    temp_audio_dir = temp_run_dir / "audio"
    try:
        temp_run_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"创建任务临时目录: {temp_run_dir}")
    except OSError as e:
         raise OSError(f"无法创建临时目录 {temp_run_dir}: {e}") from e

    # --- 1. 导出幻灯片图片 (只使用 LibreOffice) ---
    task_instance.update_state('PROCESSING', meta={'stage': STAGE_PPT_IMAGES, 'progress': 0, 'status': 'Starting slide export'})
    logger.info("--- 步骤 1: 使用 LibreOffice 导出幻灯片图片 ---")
    image_paths = None
    try:
        # 将 logger, config, task_instance 传递给导出函数
        # export_slides_with_libreoffice 内部应发送更细粒度的进度 (如导出第几页)
        image_paths = export_slides_with_libreoffice(
            pptx_filepath,
            temp_image_dir,
            logger,
            config,
            task_instance # <--- 传递任务实例
        )
        if not image_paths:
            # 导出函数内部应该记录了错误，这里直接抛出异常
            raise RuntimeError("LibreOffice 导出幻灯片图片失败或未返回任何路径。")
        logger.info(f"成功导出 {len(image_paths)} 张图片。")
        task_instance.update_state('PROCESSING', meta={'stage': STAGE_PPT_IMAGES, 'progress': 100, 'status': 'Slide export complete'})
    except Exception as e:
        logger.error(f"导出幻灯片步骤发生错误: {e}", exc_info=True)
        # export_slides_with_libreoffice 内部应已更新状态为错误
        raise RuntimeError(f"幻灯片导出失败: {e}") from e


    # --- 2. 提取备注 ---
    task_instance.update_state('PROCESSING', meta={'stage': STAGE_EXTRACT_NOTES, 'progress': 0, 'status': 'Starting notes extraction'})
    logger.info("--- 步骤 2: 提取演讲者备注 ---")
    try:
        # task_instance 参数已添加，但 extract_speaker_notes 内部暂无状态更新
        notes_list = extract_speaker_notes(pptx_filepath, logger, task_instance)
        # extract_speaker_notes 失败时会抛出异常
        task_instance.update_state('PROCESSING', meta={'stage': STAGE_EXTRACT_NOTES, 'progress': 100, 'status': 'Notes extraction complete'})
    except Exception as e:
        # logger.error(f"提取备注时出错: {e}", exc_info=True) # 内部函数已记录
        # task_instance 已在内部更新状态
        raise RuntimeError(f"提取备注失败: {e}") from e

    # --- 3. 对齐图片和备注 ---
    num_images = len(image_paths)
    num_notes = len(notes_list)
    if num_images != num_notes:
        logger.warning(f"图片数({num_images})与备注数({num_notes})不匹配，将按较小数处理。")
        min_count = min(num_images, num_notes)
        if min_count == 0:
             raise ValueError("图片或备注数量为零，无法继续处理。")
        image_paths = image_paths[:min_count]
        notes_list = notes_list[:min_count]
    elif num_images == 0:
        raise ValueError("未找到任何有效的幻灯片图片。")

    logger.info(f"将处理 {len(image_paths)} 张对齐的幻灯片/备注。")


    # --- 4. 生成音频片段 ---
    task_instance.update_state('PROCESSING', meta={'stage': STAGE_GENERATE_AUDIO, 'progress': 0, 'status': 'Starting audio generation'})
    logger.info(f"--- 步骤 4: 生成音频片段 (Edge TTS, Voice: {voice_id}) ---")
    # 确保音频目录存在
    try:
        temp_audio_dir.mkdir(parents=True, exist_ok=True)
    except OSError as e:
         raise OSError(f"无法创建音频临时目录 {temp_audio_dir}: {e}") from e

    rate_percent = config.getint('Audio', 'tts_rate_percent', fallback=100)

    # generate_audio_segments 内部处理重试和异常，并返回结果
    # 它也应发送更细粒度的进度 (如生成第几个片段)
    audio_results = generate_audio_segments(
        notes_list,
        temp_audio_dir,
        voice_id,
        rate=rate_percent,
        logger=logger,
        config=config, # 传递 config
        task_instance=task_instance # <--- 传递任务实例
    )

    if len(audio_results) != len(notes_list):
         # generate_audio_segments 应该返回与输入备注数一致的结果列表
         # 如果数量不一致，说明内部逻辑有问题
         error_msg = f"音频生成结果数量不匹配！预期 {len(notes_list)} 个，实际 {len(audio_results)} 个。"
         logger.error(error_msg)
         # task_instance 已在 generate_audio_segments 内部更新状态
         raise RuntimeError(error_msg)


    logger.info(f"音频生成过程完成。生成了 {len([p for p, d in audio_results if p])} 个有效音频片段。")
    task_instance.update_state('PROCESSING', meta={'stage': STAGE_GENERATE_AUDIO, 'progress': 100, 'status': 'Audio generation complete'})

    # --- 步骤 5: 组合结果 ---
    logger.info("--- 步骤 5: 整理处理结果 ---")
    final_data = []
    # processed_data 的长度由对齐后的图片数和备注数决定
    num_aligned_slides = len(image_paths) # 使用图片数（或对齐后的数量）作为基准

    for i in range(num_aligned_slides):
        image_path_str = image_paths[i]
        # 从 audio_results 中获取对应片段的结果
        # 即使 generate_audio_segments 有问题，结果列表长度应该与输入备注数一致
        audio_path, duration = audio_results[i]

        if not Path(image_path_str).is_file():
            logger.warning(f"最终整理时发现幻灯片 {i+1} 的图片路径无效: {image_path_str}，跳过此幻灯片。")
            continue

        slide_data = {
            'slide_number': i + 1,
            'image_path': image_path_str,
            'notes': notes_list[i] or "",
            'audio_path': audio_path, # 可能为 None
            'audio_duration': duration  # 保证是 float (0.0 或有效时长)
        }
        final_data.append(slide_data)
        logger.debug(f"  整理数据: Slide {i+1}, Img: {Path(image_path_str).name}, Audio: {Path(audio_path).name if audio_path else 'N/A'}, Dur: {duration:.3f}s")

    if not final_data:
        raise ValueError("未能整理出任何有效的幻灯片数据进行后续处理。")

    logger.info(f"成功整理了 {len(final_data)} 张幻灯片的数据。")
    # 返回处理好的数据和任务的临时目录路径给 synthesize_video_for_task
    return final_data, temp_run_dir